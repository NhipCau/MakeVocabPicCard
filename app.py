import pandas as pd
import streamlit as st
from deep_translator import GoogleTranslator
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os

# ===== UI =====
st.title("📄 絵カード作成支援ツール")

uploaded_file = st.file_uploader("ExcelまたはCSVファイルをアップロード", type=["xlsx", "csv"])

# 列指定
col_japanese = st.text_input("日本語語彙の列名または番号（A列=0）", value="0")
col_ruby = st.text_input("ルビの列名または番号（B列=1）", value="1")

# 翻訳対象言語（カンマ区切りで入力）
target_languages_str = st.text_input(
    "翻訳対象言語（カンマ区切り）  [言語コードはこちら](https://cloud.google.com/translate/docs/languages?hl=ja)",
    value="en,ne,vi,my,zh-CN,zh-TW"
)
target_languages = [lang.strip() for lang in target_languages_str.split(",") if lang.strip()]

# 位置設定（割合）
ruby_y_percent = st.number_input("ルビの縦位置（％）", value=75) / 100
word_y_percent = st.number_input("日本語語彙の縦位置（％）", value=78) / 100
translation_y_percent = st.number_input("翻訳語の縦位置（％）", value=85) / 100

# フォントサイズ設定
font_size_ruby = st.number_input("ルビのフォントサイズ", value=20)
font_size_word = st.number_input("日本語語彙のフォントサイズ", value=36)
font_size_translation = st.number_input("翻訳語のフォントサイズ", value=24)

# スライドサイズ（EMU）
SLIDE_WIDTH = 914400 * 10
SLIDE_HEIGHT = 914400 * 7.5

# ===== 関数 =====
def add_textbox(slide, text, y_percent, font_size):
    textbox = slide.shapes.add_textbox(
        left=int(SLIDE_WIDTH * 0.05),
        top=int(SLIDE_HEIGHT * y_percent),
        width=int(SLIDE_WIDTH * 0.90),
        height=int(SLIDE_HEIGHT * 0.10)
    )
    tf = textbox.text_frame
    tf.word_wrap = True   # 自動改行ON
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.alignment = PP_ALIGN.CENTER

def translate_word(word, lang):
    try:
        return GoogleTranslator(source='ja', target=lang).translate(word)
    except Exception:
        return f"[Error:{lang}]"

def create_ppt_from_vocab(df, col_japanese, col_ruby, base_filename):
    prs = Presentation()
    for _, row in df.iterrows():
        # 列の取得（番号 or 列名）
        if isinstance(col_japanese, int):
            word = str(row.iloc[col_japanese]).strip()
        else:
            word = str(row[col_japanese]).strip()

        if isinstance(col_ruby, int):
            ruby = str(row.iloc[col_ruby]).strip()
        else:
            ruby = str(row[col_ruby]).strip()

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        add_textbox(slide, ruby, ruby_y_percent, font_size_ruby)
        add_textbox(slide, word, word_y_percent, font_size_word)

        translations = [translate_word(word, lang) for lang in target_languages]
        add_textbox(slide, "   ".join(translations), translation_y_percent, font_size_translation)

    output_pptx = f"{base_filename}.pptx"
    prs.save(output_pptx)
    return output_pptx

# ===== 実行部分 =====
if uploaded_file:
    # 列指定の型変換
    try:
        col_japanese = int(col_japanese)
    except ValueError:
        pass
    try:
        col_ruby = int(col_ruby)
    except ValueError:
        pass

    if uploaded_file.name.endswith(".xlsx"):
        df = pd.read_excel(uploaded_file)
    else:
        df = pd.read_csv(uploaded_file)

    base_filename = os.path.splitext(uploaded_file.name)[0]  # 拡張子除去

    if st.button("PPTを作成"):
        ppt_path = create_ppt_from_vocab(df, col_japanese, col_ruby, base_filename)
        with open(ppt_path, "rb") as f:
            st.download_button(
                label="📥 PPTXをダウンロード",
                data=f,
                file_name=f"{base_filename}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
